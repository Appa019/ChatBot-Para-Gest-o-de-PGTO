import streamlit as st
import os
import zipfile
import tempfile
import shutil
from pathlib import Path
import time

# Imports do LlamaIndex
from llama_index.core import SimpleDirectoryReader, VectorStoreIndex, Settings, Document
from llama_index.core.node_parser import SentenceSplitter
from llama_index.llms.openai import OpenAI

# Import condicional para embeddings
try:
    from llama_index.embeddings.huggingface import HuggingFaceEmbedding
except ImportError:
    HuggingFaceEmbedding = None

# Configuração da página
st.set_page_config(
    page_title="RAG Interativo-BackOffice e Gestao",
    page_icon="📚",
    layout="wide",
    initial_sidebar_state="collapsed"
)

# CSS customizado para interface minimalista
st.markdown("""
<style>
    .main-header {
        text-align: center;
        padding: 2rem 0;
        background: linear-gradient(90deg, #667eea 0%, #764ba2 100%);
        color: white;
        border-radius: 10px;
        margin-bottom: 2rem;
    }
    .upload-section {
        background: #f8f9fa;
        padding: 2rem;
        border-radius: 10px;
        border: 2px dashed #dee2e6;
        text-align: center;
        margin-bottom: 2rem;
    }
    .chat-container {
        background: white;
        border-radius: 10px;
        padding: 1rem;
        box-shadow: 0 2px 4px rgba(0,0,0,0.1);
    }
    .success-message {
        background: #d4edda;
        color: #155724;
        padding: 1rem;
        border-radius: 5px;
        border: 1px solid #c3e6cb;
        margin: 1rem 0;
    }
    .error-message {
        background: #f8d7da;
        color: #721c24;
        padding: 1rem;
        border-radius: 5px;
        border: 1px solid #f5c6cb;
        margin: 1rem 0;
    }
    .stButton > button {
        width: 100%;
        background: linear-gradient(90deg, #667eea 0%, #764ba2 100%);
        color: white;
        border: none;
        padding: 0.5rem 1rem;
        border-radius: 5px;
        font-weight: bold;
    }
    .stTextInput > div > div > input {
        border-radius: 5px;
    }
</style>
""", unsafe_allow_html=True)

def processar_pptx(caminho_arquivo):
    """Processa arquivo PowerPoint extraindo texto de todos os slides"""
    try:
        from pptx import Presentation
        
        prs = Presentation(caminho_arquivo)
        slides_texto = []
        
        for i, slide in enumerate(prs.slides, 1):
            slide_content = f"\n=== SLIDE {i} ===\n"
            
            # Extrair texto de todas as formas no slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content += f"{shape.text}\n"
                
                # Extrair texto de tabelas se houver
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join([cell.text.strip() for cell in row.cells])
                        if row_text.strip():
                            slide_content += f"Tabela: {row_text}\n"
            
            slides_texto.append(slide_content)
        
        texto_completo = "\n".join(slides_texto)
        
        return Document(
            text=texto_completo,
            metadata={
                "file_name": os.path.basename(caminho_arquivo),
                "file_type": "pptx",
                "total_slides": len(prs.slides),
                "source": "presentation"
            }
        ), len(prs.slides)
        
    except Exception as e:
        st.error(f"Erro ao processar PPTX {os.path.basename(caminho_arquivo)}: {e}")
        return None, 0

def processar_documentos(pasta_temp):
    """Processa todos os documentos suportados na pasta temporária"""
    
    formatos_suportados = ['.pdf', '.docx', '.pptx', '.txt', '.md']
    arquivos = []
    
    # Encontrar todos os arquivos suportados
    for arquivo in os.listdir(pasta_temp):
        extensao = Path(arquivo).suffix.lower()
        if extensao in formatos_suportados:
            arquivos.append(arquivo)
    
    if not arquivos:
        return [], 0
    
    documentos_processados = []
    progress_bar = st.progress(0)
    status_text = st.empty()
    
    for i, arquivo in enumerate(arquivos):
        caminho = os.path.join(pasta_temp, arquivo)
        extensao = Path(arquivo).suffix.lower()
        
        status_text.text(f"Processando: {arquivo}")
        progress_bar.progress((i + 1) / len(arquivos))
        
        try:
            if extensao == ".pptx":
                # Processamento específico para PowerPoint
                doc, num_slides = processar_pptx(caminho)
                if doc:
                    documentos_processados.append(doc)
                    
            elif extensao in [".docx", ".pdf", ".txt", ".md"]:
                # Usar loader padrão para outros formatos
                loader = SimpleDirectoryReader(
                    input_files=[caminho],
                    required_exts=[extensao]
                )
                docs = loader.load_data()
                
                # Adicionar metadata específico
                for doc in docs:
                    doc.metadata.update({
                        "file_name": arquivo,
                        "file_type": extensao[1:],
                        "source": "document"
                    })
                
                documentos_processados.extend(docs)
                
        except Exception as e:
            st.warning(f"Erro ao processar {arquivo}: {e}")
    
    progress_bar.progress(1.0)
    status_text.text("Processamento concluído!")
    time.sleep(1)
    progress_bar.empty()
    status_text.empty()
    
    return documentos_processados, len(arquivos)

def configurar_sistema():
    """Configura o sistema RAG com os modelos"""
    try:
        # Configurar embeddings se disponível
        if HuggingFaceEmbedding:
            embed_model = HuggingFaceEmbedding(model_name="sentence-transformers/all-MiniLM-L6-v2")
            Settings.embed_model = embed_model
        
        Settings.llm = OpenAI(model="gpt-3.5-turbo", temperature=0.5)
        return True
    except Exception as e:
        st.error(f"Erro na configuração: {e}")
        return False

def criar_indice(documentos):
    """Cria o índice vetorial dos documentos"""
    try:
        with st.spinner("Criando índice vetorial..."):
            node_parser = SentenceSplitter(
                chunk_size=3000, 
                chunk_overlap=600,
                paragraph_separator="\n\n"
            )
            
            index = VectorStoreIndex.from_documents(
                documentos, 
                node_parser=node_parser
            )
            

            query_engine = index.as_query_engine(
                similarity_top_k=40,             # Mais documentos consultados
                response_mode="tree_summarize",         # Mantém o modo
                verbose=True,
                streaming=False,
                max_tokens=12000,                        # Limite maior de resposta
                temperature=0.5
            )
            
            return query_engine
    except Exception as e:
        st.error(f"Erro ao criar índice: {e}")
        return None

def modal_api_key():
    """Modal para inserir API Key"""
    if 'api_key_configured' not in st.session_state:
        st.session_state.api_key_configured = False
    
    if not st.session_state.api_key_configured:
        # Remover o CSS do modal overlay que está causando problema
        st.markdown("### 🔑 Configuração Inicial")
        st.info("Para usar o sistema, insira sua OpenAI API Key:")
        
        # Campo de input direto (sem colunas que podem causar problema)
        api_key = st.text_input(
            "OpenAI API Key", 
            type="password", 
            placeholder="sk-...",
            key="api_key_input",
            help="Sua API Key da OpenAI (começa com sk-)"
        )
        
        if st.button("🚀 Configurar Sistema", key="config_button", use_container_width=True):
            if api_key and api_key.startswith("sk-"):
                os.environ["OPENAI_API_KEY"] = api_key
                if configurar_sistema():
                    st.session_state.api_key_configured = True
                    st.success("✅ Sistema configurado com sucesso!")
                    st.rerun()
                else:
                    st.error("❌ Erro na configuração. Verifique sua API Key.")
            else:
                st.error("❌ Por favor, insira uma API Key válida (deve começar com 'sk-')")
        
        return False
    
    return True


def main():
    """Função principal da aplicação"""
    
    # Verificar e configurar API Key
    if not modal_api_key():
        return
    
    # Header principal
    st.markdown("""
    <div class="main-header">
        <h1>📚 RAG Interativo-Explicacoes BackOffice e Gestao</h1>
        <p>ChatBot</p>
    </div>
    """, unsafe_allow_html=True)
    
    # Inicializar session state
    if 'documentos_processados' not in st.session_state:
        st.session_state.documentos_processados = False
        st.session_state.query_engine = None
        st.session_state.num_documentos = 0
    
    # Seção de upload
    if not st.session_state.documentos_processados:
        st.markdown("""
        <div class="upload-section">
            <h3>📁 Upload de Documentos</h3>
            <p>Faça upload de um arquivo ZIP contendo seus documentos</p>
            <p><strong>Formatos suportados:</strong> PDF, DOCX, PPTX, TXT, MD</p>
        </div>
        """, unsafe_allow_html=True)
        
        uploaded_file = st.file_uploader(
            "Selecione um arquivo ZIP",
            type=['zip'],
            help="O ZIP deve conter os documentos que você deseja analisar"
        )
        
        if uploaded_file is not None:
            with st.spinner("Extraindo e processando documentos..."):
                # Criar pasta temporária
                with tempfile.TemporaryDirectory() as temp_dir:
                    # Salvar e extrair ZIP
                    zip_path = os.path.join(temp_dir, "upload.zip")
                    with open(zip_path, "wb") as f:
                        f.write(uploaded_file.getbuffer())
                    
                    # Extrair ZIP
                    extract_dir = os.path.join(temp_dir, "extracted")
                    with zipfile.ZipFile(zip_path, 'r') as zip_ref:
                        zip_ref.extractall(extract_dir)
                    
                    # Processar documentos
                    documentos, num_arquivos = processar_documentos(extract_dir)
                    
                    if documentos:
                        # Criar índice
                        query_engine = criar_indice(documentos)
                        
                        if query_engine:
                            st.session_state.documentos_processados = True
                            st.session_state.query_engine = query_engine
                            st.session_state.num_documentos = len(documentos)
                            
                            st.markdown(f"""
                            <div class="success-message">
                                ✅ <strong>{len(documentos)} documentos processados com sucesso!</strong>
                            </div>
                            """, unsafe_allow_html=True)
                            
                            st.rerun()
                        else:
                            st.markdown("""
                            <div class="error-message">
                                ❌ Erro ao criar índice dos documentos
                            </div>
                            """, unsafe_allow_html=True)
                    else:
                        st.markdown("""
                        <div class="error-message">
                            ❌ Nenhum documento válido encontrado no ZIP
                        </div>
                        """, unsafe_allow_html=True)
    
    # Seção de chat
    else:
        st.markdown(f"""
        <div class="success-message">
            ✅ <strong>{st.session_state.num_documentos} documentos carregados</strong> - Sistema pronto para perguntas!
        </div>
        """, unsafe_allow_html=True)
        
        # Botão para recarregar documentos
        if st.button("🔄 Carregar Novos Documentos"):
            st.session_state.documentos_processados = False
            st.session_state.query_engine = None
            st.session_state.num_documentos = 0
            st.rerun()
        
        st.markdown("---")
        
        # Interface de chat
        st.markdown("""
        <div class="chat-container">
            <h3>💬 Faça sua pergunta</h3>
        </div>
        """, unsafe_allow_html=True)
        
        # Inicializar histórico de chat
        if 'chat_history' not in st.session_state:
            st.session_state.chat_history = []
        
        # Campo de pergunta
        pergunta = st.text_input(
            "Digite sua pergunta sobre os documentos:",
            placeholder="Ex: Faça um resumo dos documentos carregados",
            key="pergunta_input"
        )
        
        col1, col2 = st.columns([3, 1])
        with col2:
            enviar = st.button("Enviar", key="enviar_pergunta")
        
        # Processar pergunta
        if enviar and pergunta.strip():
            with st.spinner("Analisando documentos..."):
                try:
                    response = st.session_state.query_engine.query(pergunta)
                    
                    st.session_state.chat_history.append({
                        "pergunta": pergunta,
                        "resposta": response.response
                    })
                    
                except Exception as e:
                    st.error(f"Erro ao processar pergunta: {e}")
        
        # Exibir histórico de chat
        if st.session_state.chat_history:
            st.markdown("---")
            st.markdown("### 📝 Histórico de Conversas")
            
            for i, chat in enumerate(reversed(st.session_state.chat_history)):
                with st.expander(f"💭 {chat['pergunta'][:50]}...", expanded=(i==0)):
                    st.markdown(f"**Pergunta:** {chat['pergunta']}")
                    st.markdown(f"**Resposta:** {chat['resposta']}")

if __name__ == "__main__":
    main()

